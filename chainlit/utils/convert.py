from PyPDF2 import PdfReader
import re
import pptx
import json

def extract_text_from_pdf(pdf_file_path):
    text = ""
    with open(pdf_file_path, "rb") as f:
        reader = PdfReader(f)
        for page in reader.pages:
            text += page.extract_text()
    return text

def file2text(file,client=None):
    print(file)
    if re.search(r'\.mp3$',file): #mp3 -> string
        audio_file = open(file, "rb")
        transcription = client.audio.transcriptions.create(
            model="whisper-1",
            file=audio_file
        )
        return transcription.text

    elif re.search(r'\.pdf$',file): #pdf -> string
        text = ""
        with open(file, "rb") as f:
            reader = PdfReader(f)
            for page in reader.pages:
                text += page.extract_text()
        return text

    elif re.search(r'\.ppt$',file) or re.search(r'\.pptx$',file): #ppt -> string
        text = ""
        presentation = pptx.Presentation(file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text
        return text

    else: #txt -> string
        try:
            with open(file, 'r', encoding='utf-8') as file:
                text = file.read()
            return text
        except FileNotFoundError:
            print("파일을 찾을 수 없습니다.")
            return None
        except Exception as e:
            print("오류 발생:", e)
            return None

def respoens2(client,model,data,query,temperature=1,max_tokens=1500):
    response = client.chat.completions.create(
        model=model,
        messages=[
            {
                "role": "user",
                "content": data
            },
            {
                "role": "assistant",
                "content": query
            }

        ],
        temperature=temperature,
        max_tokens=max_tokens,
        top_p=1,
        frequency_penalty=0,
        presence_penalty=0
    )
    return response.choices[0].message.content

def content_preprocessing2(client,file,style=None):
    model = "gpt-4o"
    #model = "gpt-4o" #한번에 데이터를 넣을 수 있음.

    pdf_text = file2text(file,client) #ppt,mp3,txt,pdf -> string으로 return

    x = len(pdf_text)
    num = x // 2
    
    pdf_1 = '<<'+pdf_text[:num]+'>>'
    pdf_2 = '<<'+pdf_text[num:]+'>>'
    full_pdf = '<<'+pdf_text+'>>'
    

    command1 = "이를 참고하여 << >>안에 있는 내용을 바탕으로 워드 한페이지 분량의 교재를 만들어라. 내용은 정확하고 자세해야 한다."

    chunk1 = respoens2(client,model,pdf_1,command1)
    chunk2 = respoens2(client,model,pdf_2,command1)

    enter_str = "\n\n"

    re_str = full_pdf


    return re_str

def extract_from_json_file(file_path):

    with open(file_path, 'r', encoding='utf-8') as file:
        data = json.load(file)

    questions = data.get('question', [])
    answers = data.get('answer', [])

    return questions, answers

def validate_answer(client,Q_list,A_list,context=False,upgrade_question=False):

    context_str = ''
    context_str2 = ''
    if context:
        context_str += f'''### Context
{context}
'''
        context_str2 += 'Answer using the information in the context as much as possible'

    #context가 있으면 위에 코드도 같이 실행됨.

    len_quiz = len(Q_list)
    for i in range(len_quiz):
        command_str = f'''
### Question:
{Q_list[i]}

### Answer:
{A_list[i]}
'''
        command_str2 = f'''
### Instruction:
Is the answer to the question correct?
If it is not correct, Give me the correct answer just in a similar format to <{A_list[i]}>
or If it's accurate, just print out the answer as it is.
'''

        answer = respoens2(client, 'gpt-4o', command_str+context_str, command_str2+context_str2, 0)

        if answer.lower() != A_list[i].lower():
            A_list[i] = answer

    return Q_list, A_list


def save_txt(Q_list, A_list, file_path):

    with open(file_path, 'w', encoding='utf-8') as file:

        Q_str = ''
        A_str = ''
        for i in range(len(Q_list)):
            Q_str += f"{Q_list[i]}"

        for j in range(len(A_list)):
            A_str += f"{A_list[j]}"

        file.write("Question:"+Q_str+"\n")
        file.write("Answer:"+A_str)

def save_qa_to_json(Q_list, A_list, file_path):

    data = {
        "questions": Q_list,
        "answers": A_list
    }


    with open(file_path, 'w', encoding='utf-8') as file:
        json.dump(data, file, indent=4)

def validate_and_save(client,json_file_path,txt_file_path=False):
    '''
    JSON file -> 질문,답변 가져오기 -> 검증 (-> txt파일에 저장) -> JSON file

    '''
    Q_list, A_list = extract_from_json_file(json_file_path)

    Q, A = validate_answer(client,Q_list,A_list) #context 필요시 작성.

    if txt_file_path:
        save_txt(Q,A,txt_file_path)

    save_qa_to_json(Q,A)
