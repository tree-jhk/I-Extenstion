from PyPDF2 import PdfReader
import re
import pptx
import json
from PIL import Image
import pytesseract
from pdf2image import convert_from_path
import os
OUTPUT_PATH = r'chainlit/utils/tempImage'

'''
POPPLER_PATH와 pytesseract.pytesseract.tesseract_cmd는 bin의 절대경로, tesseract.exe의 절대경로를 가져와야합니다.
ctrl+shift+c를 통해서 그 파일들의 절대경로를 가져와야함;;;
'''

POPPLER_PATH = r'C:\Users\qkrck\PycharmProjects\PBPProject\I-Quiz\chainlit\utils\poppler-24.02.0\Library\bin'
pytesseract.pytesseract.tesseract_cmd = r'C:\Users\qkrck\PycharmProjects\PBPProject\I-Quiz\chainlit\utils\새 폴더\tesseract.exe'

NAMEOFFILE="qa_file.json"
MAXTOKEN = 4096


def pdf_to_images(pdf_path, output_folder):
    pages = convert_from_path(pdf_path, poppler_path=POPPLER_PATH)

    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    image_paths = []
    for i, page in enumerate(pages):
        image_path = os.path.join(output_folder, f'page_{i + 1}.jpg')
        page.save(image_path, 'JPEG')
        image_paths.append(image_path)


    return image_paths
def extract_text_from_image(image_path):
    image = Image.open(image_path)

    text = pytesseract.image_to_string(image, lang='eng+kor')

    return text

def extract_text_from_pdf(c,pdf_file_path):
    text1 = ""
    with open(pdf_file_path, "rb") as f:
        reader = PdfReader(f)
        for page in reader.pages:
            text1 += page.extract_text()


    x = pdf_to_images(pdf_file_path,OUTPUT_PATH)

    if len(text1) < 2000:
        text2 = ""
        for pt in x:
            text2 += extract_text_from_image(pt)

        return text2

    return text1

def file2text(file,client=None):

    if re.search(r'\.mp3$',file): #mp3 -> string
        audio_file = open(file, "rb")
        transcription = client.audio.transcriptions.create(
            model="whisper-1",
            file=audio_file
        )
        return transcription.text

    elif re.search(r'\.pdf$',file): #pdf -> string 2222222222222222222222222222222222222222222222222222222222222222222222222222222222222222222222222222222
        text = extract_text_from_pdf(client,file)
        return text

    elif re.search(r'\.png$',file) or re.search(r'\.jpg$',file) :
        text = extract_text_from_image(file)
        return text

    elif re.search(r'\.ppt$',file) or re.search(r'\.pptx$',file): #ppt -> string
        text = ""
        presentation = pptx.Presentation(file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text
        return text

    else: #txt -> string
        try:
            with open(file, 'r', encoding='utf-8') as file:
                text = file.read()
            return text
        except FileNotFoundError:
            print("파일을 찾을 수 없습니다.")
            return None
        except Exception as e:
            print("오류 발생:", e)
            return None

def respoens2(client,model,data,query,temperature=1,max_tokens=MAXTOKEN):

    response = client.chat.completions.create(
        model=model,
        messages=[
            {
                "role": "user",
                "content": data
            },
            {
                "role": "assistant",
                "content": query
            }

        ],
        temperature=temperature,
        max_tokens=max_tokens,
        top_p=1,
        frequency_penalty=0,
        presence_penalty=0
    )

    return response.choices[0].message.content

def content_preprocessing2(client,file,style=None):
    model = "gpt-4o"
    #model = "gpt-4o" #한번에 데이터를 넣을 수 있음.

    pdf_text = file2text(file,client) #ppt,mp3,txt,pdf -> string으로 return

    x = len(pdf_text)
    num = x // 2
    
    pdf_1 = '<<'+pdf_text[:num]+'>>'
    pdf_2 = '<<'+pdf_text[num:]+'>>'
    full_pdf = '<<'+pdf_text+'>>'
    

    command1 = "이를 참고하여 << >>안에 있는 내용을 바탕으로 워드 한페이지 분량의 교재를 만들어라. 내용은 정확하고 자세해야 한다."

    chunk1 = respoens2(client,model,pdf_1,command1)
    chunk2 = respoens2(client,model,pdf_2,command1)

    enter_str = "\n\n"

    re_str = full_pdf


    return re_str



def validate_answer(client,Q_list,A_list,context=False,upgrade_question=False):

    context_str = ''
    context_str2 = ''
    if context:
        context_str += f'''### Context
{context}
'''
        context_str2 += 'Answer using the information in the context as much as possible'

    #context가 있으면 위에 코드도 같이 실행됨.

    len_quiz = len(Q_list)
    for i in range(len_quiz):
        command_str = f'''
### Question:
{Q_list[i]}

### Answer:
{A_list[i]}
'''
        command_str2 = f'''
### Instruction:
Is the answer to the question correct?
If it is not correct, Give me the correct answer just in a similar format to <{A_list[i]}>
or If it's accurate, just print out the answer as it is.
'''

        answer = respoens2(client, 'gpt-4o', command_str+context_str, command_str2+context_str2, 0)

        if answer.lower() != A_list[i].lower():
            A_list[i] = answer

    return Q_list, A_list


#txt file로 변환해주는 기능.
def save_txt(Q_list, A_list, file_path):

    with open(file_path, 'w', encoding='utf-8') as file:
        Q_str = '\n'.join(Q_list)
        A_str = '\n'.join(A_list)

        file.write("Question:\n"+Q_str+"\n\n")
        file.write("Answer:\n"+A_str)


'''
JSON FILE 처리하는 함수가 있는 부분.
'''

def save_qa_to_json(questions, answers, filename=NAMEOFFILE):
    if len(questions) != len(answers):
        raise ValueError("The number of questions and answers must be the same.")

    qa_pairs = [{"question": q, "answer": a} for q, a in zip(questions, answers)]

    with open(filename, 'w', encoding='utf-8') as f:
        json.dump(qa_pairs, f, ensure_ascii=False, indent=4)

def extract_from_json_file(filename=NAMEOFFILE):
    with open(filename, 'r', encoding='utf-8') as f:
        data = json.load(f)

    questions = [item['question'] for item in data]
    answers = [item['answer'] for item in data]

    return questions, answers


'''
검증하는 함수 통합이기는 하지만 수정 필요할 것으로 예상!
'''
def validate_and_save(client,json_file_path,txt_file_path=False):
    '''
    JSON file -> 질문,답변 가져오기 -> 검증 (-> txt파일에 저장) -> JSON file
    '''
    Q_list, A_list = extract_from_json_file(json_file_path)

    Q, A = validate_answer(client,Q_list,A_list) #context 필요시 작성.

    if txt_file_path:
        save_txt(Q,A,txt_file_path)

    save_qa_to_json(Q,A)


def check_a(client,q,user_a):
    model = "gpt-4o"
    command1 = f"질문:{q}"
    command2 = f"질문에 대한 답이 {user_a}이 맞는가? 맞는지 틀린지를 먼저 기술하고, 부연설명은 짧게 해라."

    an = respoens2(client, model, command1, command2)

    return str(an)


def safe_message(client, m, num):
    model = "gpt-4o"
    m = str(m)
    x = len(m) // num
    lists = [m[i:i + x] for i in range(0, len(m), x)]

    output = ''
    for j in range(len(lists)):
        chunk = respoens2(client,model,' ',lists[j])
        output += chunk

    return output
